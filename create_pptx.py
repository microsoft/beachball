#!/usr/bin/env python3
"""Generate presentation.pptx by parsing presentation-outline.md with markdown-it-py."""

from pathlib import Path
from markdown_it import MarkdownIt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import re

MD_PATH = Path(__file__).parent / "presentation-outline.md"
OUT_PATH = Path(__file__).parent / "presentation.pptx"


# --- Markdown parsing ---


def inline_text(token):
    """Extract plain text from an inline token, stripping markdown formatting."""
    if token.children:
        return "".join(c.content for c in token.children if c.type in ("text", "code_inline", "softbreak"))
    return token.content


def parse_markdown(path):
    """Parse the markdown file into a list of slide dicts."""
    md = MarkdownIt("commonmark").enable("table")
    tokens = md.parse(path.read_text())

    slides = []
    current = None
    i = 0

    while i < len(tokens):
        t = tokens[i]

        # --- H1: title slide ---
        if t.type == "heading_open" and t.tag == "h1":
            title = inline_text(tokens[i + 1])
            current = {"type": "title", "title": title, "subtitle": ""}
            slides.append(current)
            i += 3  # skip heading_open, inline, heading_close
            continue

        # Subtitle (italic paragraph right after h1)
        if (current and current["type"] == "title" and not current["subtitle"]
                and t.type == "paragraph_open"):
            text = inline_text(tokens[i + 1])
            # Check if the inline's children are just an em wrapper
            children = tokens[i + 1].children or []
            if any(c.type == "em_open" for c in children):
                current["subtitle"] = text
                i += 3  # paragraph_open, inline, paragraph_close
                continue

        # --- H2: new content slide ---
        if t.type == "heading_open" and t.tag == "h2":
            title = inline_text(tokens[i + 1])
            current = {"type": "content", "title": title, "body": [], "notes": None}
            slides.append(current)
            i += 3
            continue

        if current is None or current["type"] == "title":
            i += 1
            continue

        # --- HTML comments: skip ---
        if t.type == "html_block":
            i += 1
            continue

        # --- Slide notes: (Slide notes: ...) ---
        if t.type == "paragraph_open":
            text = inline_text(tokens[i + 1])
            if text.startswith("(Slide notes:"):
                notes = re.sub(r"^\(Slide notes:\s*", "", text)
                notes = re.sub(r"\)$", "", notes)
                current["notes"] = notes
                i += 3
                continue

        # --- Code fence ---
        if t.type == "fence":
            current["body"].append({"kind": "code", "lines": t.content.rstrip("\n").split("\n"), "info": t.info})
            i += 1
            continue

        # --- Table ---
        if t.type == "table_open":
            headers = []
            rows = []
            i += 1
            in_thead = False
            in_tbody = False
            current_row = []
            while i < len(tokens) and tokens[i].type != "table_close":
                tt = tokens[i]
                if tt.type == "thead_open":
                    in_thead = True
                elif tt.type == "thead_close":
                    in_thead = False
                elif tt.type == "tbody_open":
                    in_tbody = True
                elif tt.type == "tbody_close":
                    in_tbody = False
                elif tt.type == "tr_open":
                    current_row = []
                elif tt.type == "tr_close":
                    if in_thead:
                        headers = current_row
                    else:
                        rows.append(current_row)
                elif tt.type == "inline":
                    current_row.append(inline_text(tt))
                i += 1
            i += 1  # skip table_close
            current["body"].append({"kind": "table", "headers": headers, "rows": rows})
            continue

        # --- Bullet list ---
        if t.type == "bullet_list_open":
            items = []
            i += 1
            list_base_level = t.level
            while i < len(tokens) and tokens[i].type != "bullet_list_close":
                tt = tokens[i]
                if tt.type == "list_item_open":
                    indent = max(0, (tt.level - list_base_level) // 2)
                    # Find the inline content
                    for j in range(i + 1, min(i + 6, len(tokens))):
                        if tokens[j].type == "inline":
                            items.append({"text": inline_text(tokens[j]), "indent": indent})
                            break
                    # Handle nested lists: don't skip them, let the loop process them
                elif tt.type == "bullet_list_open":
                    # Nested list — walk through it extracting items
                    nested_base = tt.level
                    i += 1
                    while i < len(tokens) and tokens[i].type != "bullet_list_close":
                        nt = tokens[i]
                        if nt.type == "list_item_open":
                            indent = max(0, (nt.level - list_base_level) // 2)
                            for j in range(i + 1, min(i + 6, len(tokens))):
                                if tokens[j].type == "inline":
                                    items.append({"text": inline_text(tokens[j]), "indent": indent})
                                    break
                        i += 1
                    # i now points at bullet_list_close for nested list
                i += 1
            i += 1  # skip bullet_list_close
            current["body"].append({"kind": "bullets", "items": items})
            continue

        # --- Plain paragraph ---
        if t.type == "paragraph_open":
            text = inline_text(tokens[i + 1])
            if text.strip():
                current["body"].append({"kind": "text", "text": text})
            i += 3
            continue

        # --- Horizontal rule (ignored) ---
        if t.type == "hr":
            i += 1
            continue

        i += 1

    return slides


# --- PowerPoint generation ---

prs = Presentation()

TITLE_LAYOUT = prs.slide_layouts[0]
CONTENT_LAYOUT = prs.slide_layouts[1]
TITLE_ONLY_LAYOUT = prs.slide_layouts[5]


def clean(text):
    """Strip markdown formatting for plain-text placeholders."""
    text = text.replace("`", "")
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"_(.+?)_", r"\1", text)
    return text


def remove_bullet(p):
    """Remove the bullet from a paragraph."""
    pPr = p._p.get_or_add_pPr()
    for child in list(pPr):
        if child.tag.endswith(("buChar", "buAutoNum", "buNone", "buFont")):
            pPr.remove(child)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def render_slide(s):
    if s["type"] == "title":
        slide = prs.slides.add_slide(TITLE_LAYOUT)
        slide.placeholders[0].text = s["title"]
        slide.placeholders[1].text = s.get("subtitle", "")
        return

    # Determine slide type based on body content
    has_table = any(b["kind"] == "table" for b in s["body"])
    has_code = any(b["kind"] == "code" for b in s["body"])

    if has_table:
        slide = prs.slides.add_slide(TITLE_ONLY_LAYOUT)
        slide.placeholders[0].text = s["title"]

        for block in s["body"]:
            if block["kind"] == "table":
                headers = block["headers"]
                rows = block["rows"]
                num_rows = len(rows) + 1
                num_cols = len(headers)
                table_shape = slide.shapes.add_table(
                    num_rows, num_cols,
                    Inches(0.5), Inches(1.6),
                    Inches(9.0), Inches(0.4 * num_rows),
                )
                table = table_shape.table
                for ci, h in enumerate(headers):
                    cell = table.cell(0, ci)
                    cell.text = clean(h)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(14)
                        p.font.bold = True
                for ri, row in enumerate(rows):
                    for ci, val in enumerate(row):
                        cell = table.cell(ri + 1, ci)
                        cell.text = clean(val)
                        for p in cell.text_frame.paragraphs:
                            p.font.size = Pt(13)

    elif has_code:
        slide = prs.slides.add_slide(CONTENT_LAYOUT)
        slide.placeholders[0].text = s["title"]
        tf = slide.placeholders[1].text_frame
        tf.clear()
        first = True

        for block in s["body"]:
            if block["kind"] == "code":
                for code_line in block["lines"]:
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    p.text = code_line
                    p.font.size = Pt(14)
                    p.font.name = "Courier New"
                    p.space_after = Pt(2)
                    remove_bullet(p)
            elif block["kind"] == "text":
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                    p.text = ""
                    remove_bullet(p)
                    p = tf.add_paragraph()
                p.text = clean(block["text"])
                p.font.size = Pt(18)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                remove_bullet(p)

    else:
        # Content slide with bullets/text
        slide = prs.slides.add_slide(CONTENT_LAYOUT)
        slide.placeholders[0].text = s["title"]
        tf = slide.placeholders[1].text_frame
        tf.clear()
        first = True

        for block in s["body"]:
            if block["kind"] == "bullets":
                for item in block["items"]:
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    p.text = clean(item["text"])
                    p.level = item["indent"]
                    p.space_after = Pt(4)
                    p.font.size = Pt(18) if item["indent"] == 0 else Pt(16)
            elif block["kind"] == "text":
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                p.text = clean(block["text"])
                p.space_after = Pt(4)
                p.font.size = Pt(18)

    # Speaker notes
    if s.get("notes"):
        slide.notes_slide.notes_text_frame.text = s["notes"]


# --- Main ---

slides = parse_markdown(MD_PATH)
for s in slides:
    render_slide(s)
prs.save(str(OUT_PATH))
print(f"Saved {len(prs.slides)} slides to {OUT_PATH}")
